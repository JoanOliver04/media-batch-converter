"""Readers and writers for DOCX, XLSX, CSV and PPTX."""

from __future__ import annotations

import csv
from pathlib import Path

from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt

from documents.errors import DocumentError
from documents.model import Block, BlockKind, DocumentModel
from documents.security import MAX_BLOCKS, MAX_TABLE_CELLS
from documents.textio import read_text_file
from i18n import t


def read_docx(path: Path) -> DocumentModel:
    document = Document(str(path))
    blocks: list[Block] = []
    cells = 0
    title: str | None = None
    for item in _iter_docx_items(document):
        _ensure_block_budget(len(blocks))
        if isinstance(item, Paragraph):
            text = item.text.strip()
            if not text:
                continue
            style = (item.style.name if item.style is not None else "") or ""
            if style.startswith("Title") and title is None:
                title = text
                continue
            if style.startswith("Heading"):
                digits = "".join(
                    character for character in style if character.isdigit()
                )
                level = int(digits) if digits else 1
                blocks.append(Block(BlockKind.HEADING, text, level=min(level, 6)))
                continue
            if style.startswith("List"):
                blocks.append(Block(BlockKind.LIST_ITEM, text))
                continue
            blocks.append(Block(BlockKind.PARAGRAPH, text))
            continue
        rows = tuple(
            tuple(cell.text.strip() for cell in row.cells) for row in item.rows
        )
        cells += sum(len(row) for row in rows)
        if cells > MAX_TABLE_CELLS:
            raise DocumentError(t("document.too_many_cells", limit=MAX_TABLE_CELLS))
        if rows:
            blocks.append(Block(BlockKind.TABLE, rows=rows))
    return DocumentModel(title, tuple(blocks), source_format="DOCX")


def write_docx(model: DocumentModel, path: Path) -> None:
    document = Document()
    if model.title:
        document.add_heading(model.title, level=0)
    for block in model.blocks:
        if block.kind is BlockKind.HEADING:
            document.add_heading(block.text, level=min(max(block.level, 1), 9))
        elif block.kind is BlockKind.PARAGRAPH:
            document.add_paragraph(block.text)
        elif block.kind is BlockKind.LIST_ITEM:
            style = "List Number" if block.ordered else "List Bullet"
            document.add_paragraph(block.text, style=style)
        elif block.kind is BlockKind.CODE:
            paragraph = document.add_paragraph(block.text)
            for run in paragraph.runs:
                run.font.name = "Consolas"
        elif block.kind is BlockKind.TABLE:
            _write_docx_table(document, block.rows)
        elif block.kind is BlockKind.PAGE_BREAK:
            document.add_page_break()
    if not document.paragraphs and not document.tables:
        document.add_paragraph("")
    document.save(str(path))


def read_xlsx(path: Path) -> DocumentModel:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    try:
        return _workbook_to_model(workbook, "XLSX")
    finally:
        workbook.close()


def write_xlsx(model: DocumentModel, path: Path) -> None:
    workbook = Workbook()
    default = workbook.active
    tables = [block for block in model.blocks if block.kind is BlockKind.TABLE]
    if tables:
        workbook.remove(default)
        for index, block in enumerate(tables, 1):
            title = _sheet_title(model, index)
            sheet = workbook.create_sheet(title)
            for row in block.rows:
                sheet.append(list(row))
    else:
        default.title = _sheet_title(model, 1)
        if model.title:
            default.append([model.title])
        for block in model.blocks:
            if block.text:
                default.append([block.text])
        if default.max_row == 1 and default["A1"].value is None:
            default.append([""])
    workbook.save(str(path))


def read_csv(path: Path) -> DocumentModel:
    text = read_text_file(path)
    rows = tuple(tuple(row) for row in csv.reader(text.splitlines()))
    cells = sum(len(row) for row in rows)
    if cells > MAX_TABLE_CELLS:
        raise DocumentError(t("document.too_many_cells", limit=MAX_TABLE_CELLS))
    blocks = (Block(BlockKind.TABLE, rows=rows),) if rows else ()
    return DocumentModel(None, blocks, source_format="CSV")


def write_csv(model: DocumentModel, path: Path) -> None:
    table = next(
        (block.rows for block in model.blocks if block.kind is BlockKind.TABLE), None
    )
    with path.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.writer(handle)
        if table:
            writer.writerows(table)
        else:
            if model.title:
                writer.writerow([model.title])
            for block in model.blocks:
                if block.text:
                    writer.writerow([block.text])


def read_pptx(path: Path) -> DocumentModel:
    presentation = Presentation(str(path))
    blocks: list[Block] = []
    for index, slide in enumerate(presentation.slides, 1):
        _ensure_block_budget(len(blocks))
        if index > 1:
            blocks.append(Block(BlockKind.PAGE_BREAK))
        title = _slide_title(slide)
        blocks.append(
            Block(
                BlockKind.HEADING,
                title or t("document.slide_fallback", number=index),
                level=2,
            )
        )
        for text in _slide_texts(slide, skip_title=True):
            _ensure_block_budget(len(blocks))
            blocks.append(Block(BlockKind.PARAGRAPH, text))
        notes = _slide_notes(slide)
        if notes:
            blocks.append(Block(BlockKind.PARAGRAPH, notes))
    return DocumentModel(
        None, tuple(blocks), page_count=len(presentation.slides), source_format="PPTX"
    )


def write_pptx(model: DocumentModel, path: Path) -> None:
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    title_layout = presentation.slide_layouts[0]
    current = None
    body_top = Inches(1.6)

    def new_slide(heading: str) -> None:
        nonlocal current, body_top
        current = presentation.slides.add_slide(title_layout if heading else blank)
        if heading and current.shapes.title is not None:
            current.shapes.title.text = heading
        body_top = Inches(1.6)

    if model.title:
        opening = presentation.slides.add_slide(title_layout)
        if opening.shapes.title is not None:
            opening.shapes.title.text = model.title
        current = opening
    for block in model.blocks:
        if block.kind is BlockKind.HEADING or block.kind is BlockKind.PAGE_BREAK:
            new_slide(block.text if block.kind is BlockKind.HEADING else "")
            continue
        if current is None:
            new_slide(model.title or "")
        assert current is not None
        if block.kind is BlockKind.TABLE and block.rows:
            _write_pptx_table(current, block.rows)
            continue
        text = block.text.strip()
        if not text:
            continue
        box = current.shapes.add_textbox(
            Inches(0.7), body_top, Inches(8.6), Inches(0.6)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.text = f"• {text}" if block.kind is BlockKind.LIST_ITEM else text
        frame.paragraphs[0].font.size = Pt(18)
        body_top += Inches(0.45)
    if not presentation.slides:
        presentation.slides.add_slide(blank)
    presentation.save(str(path))


def _iter_docx_items(document: Document):
    body = document.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, document)
        elif child.tag == qn("w:tbl"):
            yield Table(child, document)


def _write_docx_table(document: Document, rows: tuple[tuple[str, ...], ...]) -> None:
    if not rows:
        return
    width = max(len(row) for row in rows)
    table = document.add_table(rows=len(rows), cols=width)
    table.style = "Table Grid"
    for row_index, row in enumerate(rows):
        for column, value in enumerate(row):
            table.cell(row_index, column).text = value
    if table.rows:
        for cell in table.rows[0].cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT


def _workbook_to_model(workbook, source_format: str) -> DocumentModel:
    blocks: list[Block] = []
    cells = 0
    for sheet in workbook.worksheets:
        _ensure_block_budget(len(blocks))
        blocks.append(Block(BlockKind.HEADING, sheet.title, level=2))
        rows: list[tuple[str, ...]] = []
        for row in sheet.iter_rows(values_only=True):
            values = tuple("" if cell is None else str(cell) for cell in row)
            if any(value.strip() for value in values):
                rows.append(values)
                cells += len(values)
                if cells > MAX_TABLE_CELLS:
                    raise DocumentError(
                        t("document.too_many_cells", limit=MAX_TABLE_CELLS)
                    )
        if rows:
            blocks.append(Block(BlockKind.TABLE, rows=tuple(rows)))
    return DocumentModel(None, tuple(blocks), source_format=source_format)


def _sheet_title(model: DocumentModel, index: int) -> str:
    base = (model.title or f"Sheet{index}").strip() or f"Sheet{index}"
    cleaned = "".join(
        character if character not in r"[]:*?/\\" else " " for character in base
    )
    return cleaned[:31] or f"Sheet{index}"


def _slide_title(slide) -> str:
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip()
    return ""


def _slide_texts(slide, skip_title: bool) -> list[str]:
    texts: list[str] = []
    title_shape = slide.shapes.title
    for shape in slide.shapes:
        if skip_title and shape == title_shape:
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text:
            texts.append(text)
    return texts


def _slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def _write_pptx_table(slide, rows: tuple[tuple[str, ...], ...]) -> None:
    width = max(len(row) for row in rows)
    table_shape = slide.shapes.add_table(
        len(rows), width, Inches(0.6), Inches(1.8), Inches(8.8), Inches(0.4 * len(rows))
    )
    table = table_shape.table
    for row_index, row in enumerate(rows):
        for column, value in enumerate(row):
            table.cell(row_index, column).text = value


def _ensure_block_budget(count: int) -> None:
    if count >= MAX_BLOCKS:
        raise DocumentError(t("document.too_many_blocks", limit=MAX_BLOCKS))
